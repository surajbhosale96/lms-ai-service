import io
import requests
from pptx import Presentation


def extract_pptx(url: str) -> str:
    response = requests.get(url, timeout=60)
    response.raise_for_status()

    prs = Presentation(io.BytesIO(response.content))
    texts = []

    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append(" ".join(slide_texts))

    return "\n\n".join(texts)
